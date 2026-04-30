# -*- coding: utf-8 -*-
"""
Omega CivicFlow v4 - Accurate Portfolio Presentation Generator
Based on actual codebase analysis. No exaggeration.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# -- Constants --
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

BLACK = RGBColor(0x05, 0x05, 0x05)
GOLD = RGBColor(0xC0, 0xA0, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
DIM_GRAY = RGBColor(0x66, 0x66, 0x66)
CARD_BG = RGBColor(0x14, 0x14, 0x14)
CARD_BORDER = RGBColor(0x33, 0x33, 0x33)

MARGIN_L = Emu(731520)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_L

FONT_KR = "Malgun Gothic"
FONT_EN = "Segoe UI"


def set_slide_bg(slide, color=BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=14,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name=FONT_KR, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=11,
                  color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                  bullet=False, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        prefix = "  " if bullet and line and not line.startswith(" ") else ""
        p.text = prefix + line
        p.alignment = alignment
        p.space_after = Pt(2)
        p.space_before = Pt(1)
        p.line_spacing = Pt(font_size * line_spacing)
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = FONT_KR
    return txBox


def add_card(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_divider(slide, left, top, width, color=GOLD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Emu(25400)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_header(slide, title, subtitle):
    add_text(slide, MARGIN_L, Emu(274320), CONTENT_W, Emu(457200),
             title, font_size=26, color=WHITE, bold=True)
    add_text(slide, MARGIN_L, Emu(731520), CONTENT_W, Emu(320040),
             subtitle, font_size=13, color=LIGHT_GRAY)
    add_divider(slide, MARGIN_L, Emu(1097280), Emu(2286000))


# =====================================================
# SLIDES
# =====================================================

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, Emu(0), Emu(1371600), SLIDE_W, Emu(822960),
             "\u03A9", font_size=72, color=GOLD, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, Emu(0), Emu(2377440), SLIDE_W, Emu(640080),
             "OMEGA CIVICFLOW v4", font_size=38, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_divider(slide, Emu(4114800), Emu(3200400), Emu(3962095))
    add_text(slide, Emu(0), Emu(3383280), SLIDE_W, Emu(457200),
             "DART \uae08\uc735 \uacf5\uc2dc \ubb38\uc11c \uc790\ub3d9 \ubd84\uc11d \ud50c\ub7ab\ud3fc",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text(slide, Emu(0), Emu(3840480), SLIDE_W, Emu(365760),
             "OCR x Dual-LLM x RAG x Supervisor Audit",
             font_size=14, color=MID_GRAY,
             alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_text(slide, Emu(0), Emu(4572000), SLIDE_W, Emu(457200),
             "React 18  |  FastAPI  |  EXAONE 3.5  |  Gemini 2.5 Pro/Flash  |  BGE-M3  |  ChromaDB",
             font_size=11, color=DIM_GRAY,
             alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_text(slide, Emu(0), Emu(5303520), SLIDE_W, Emu(365760),
             "Solo Full-Stack  |  2026.02 - 04  |  57,596 LoC (Python + JSX)",
             font_size=10, color=DIM_GRAY,
             alignment=PP_ALIGN.CENTER, font_name=FONT_EN)


def slide_02_toc(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "\ubaa9\ucc28", "Table of Contents")

    items = [
        ("01", "\ubb38\uc81c \uc815\uc758", "DART \uacf5\uc2dc \ubb38\uc11c \uc218\uc791\uc5c5 \ubd84\uc11d\uc758 \ud55c\uacc4"),
        ("02", "\uc194\ub8e8\uc158 -- 5\ub2e8\uacc4 \ud30c\uc774\ud504\ub77c\uc778", "\uc5c5\ub85c\ub4dc\ubd80\ud130 AI \ucc57\ubd07\uae4c\uc9c0 \uc804\uc790\ub3d9"),
        ("03", "\ud50c\ub7ab\ud3fc UI", "Monochrome Authority \ud14c\ub9c8"),
        ("04", "\ud575\uc2ec \ucc28\ubcc4\uc810 & \uc131\uacfc", "\uc65c Omega CivicFlow\uc778\uac00?"),
        ("05", "\ubc31\uc5d4\ub4dc \uc11c\ube44\uc2a4 \ub808\uc774\uc5b4", "27\uac1c \uc804\ubb38 \uc11c\ube44\uc2a4 + 4 Agent"),
        ("06", "OCR \ub0b4\ubd80 \uba54\ucee4\ub2c8\uc998", "EasyOCR + OpenCV 8\ub2e8\uacc4 \uc804\ucc98\ub9ac"),
        ("07", "LLM \ubd84\uc11d \uc5d4\uc9c4", "EXAONE 3.5 7.8B -- Ollama \ub85c\uceec"),
        ("08", "Insight + Supervisor", "Gemini 2.5 Pro + Flash \ub3c5\ub9bd \uac10\uc0ac"),
        ("09", "RAG \uac80\uc0c9 \uc2dc\uc2a4\ud15c", "BGE-M3 + \ud558\uc774\ube0c\ub9ac\ub4dc \uac80\uc0c9"),
        ("10", "\ud504\ub860\ud2b8\uc5d4\ub4dc \uad6c\ud604", "14 Pages + 4 Components"),
        ("11", "\ub370\uc774\ud130\ubca0\uc774\uc2a4 \uc2a4\ud0a4\ub9c8", "SQLite 11 Tables + ChromaDB"),
        ("12", "\ubcf4\uc548 & \uc778\uc99d", "JWT + RBAC + \ub370\uc774\ud130 \ubcf4\ud638"),
        ("13", "\uae30\uc220 \uc2a4\ud0dd & \uc544\ud0a4\ud14d\ucc98", "\uc120\ud0dd \uadfc\uac70 \ud3ec\ud568"),
    ]

    y = Emu(1371600)
    for num, title, desc in items:
        add_text(slide, MARGIN_L, y, Emu(365760), Emu(228600),
                 num, font_size=13, color=GOLD, bold=True, font_name=FONT_EN)
        add_text(slide, Emu(1188720), y, Emu(2560320), Emu(228600),
                 title, font_size=12, color=WHITE, bold=True)
        add_text(slide, Emu(3840480), y, Emu(6400800), Emu(228600),
                 desc, font_size=10, color=MID_GRAY)
        y += Emu(365760)


def slide_03_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "01  \ubb38\uc81c \uc815\uc758",
               "DART \uacf5\uc2dc \ubb38\uc11c \uc218\uc791\uc5c5 \ubd84\uc11d\uc758 \ud55c\uacc4")

    card_w = Emu(3383280)
    card_h = Emu(2560320)
    gap = Emu(274320)
    y_start = Emu(1371600)

    # Card 1
    x1 = MARGIN_L
    add_card(slide, x1, y_start, card_w, card_h)
    add_text(slide, x1 + Emu(182880), y_start + Emu(137160), card_w - Emu(365760), Emu(320040),
             "[TIME] \uc2dc\uac04 / \ube44\uc6a9", font_size=15, color=GOLD, bold=True)
    add_multiline(slide, x1 + Emu(182880), y_start + Emu(502920),
                  card_w - Emu(365760), card_h - Emu(640080),
                  ["\ubb38\uc11c 1\uac74 \ubd84\uc11d \ud3c9\uade0 15~30\ubd84 \uc18c\uc694",
                   "\ub300\ub7c9 \ubb38\uc11c \ucc98\ub9ac \uc2dc \uc2ec\uac01\ud55c \ubcd1\ubaa9",
                   "\ubd84\uc11d\uad00 1\uc778 \ud558\ub8e8 50\uac74 \ubbf8\ub9cc \ucc98\ub9ac",
                   "\uc57c\uac04/\uc8fc\ub9d0 \ucc98\ub9ac \ubd88\uac00 (\uc778\ub825 \uc758\uc874)"],
                  font_size=11, color=LIGHT_GRAY, bullet=True)

    # Card 2
    x2 = x1 + card_w + gap
    add_card(slide, x2, y_start, card_w, card_h)
    add_text(slide, x2 + Emu(182880), y_start + Emu(137160), card_w - Emu(365760), Emu(320040),
             "[QUALITY] \ud488\uc9c8 / \uc77c\uad00\uc131", font_size=15, color=GOLD, bold=True)
    add_multiline(slide, x2 + Emu(182880), y_start + Emu(502920),
                  card_w - Emu(365760), card_h - Emu(640080),
                  ["\ubd84\uc11d\uad00\ubcc4 \ud310\ub2e8 \uae30\uc900 \ubd88\uc77c\uce58",
                   "\uae08\uc735 \uc6a9\uc5b4 \ud574\uc11d \ud3b8\ucc28 \ubc1c\uc0dd",
                   "\uce74\ud14c\uace0\ub9ac \uc624\ubd84\ub958 -> \uc758\uc0ac\uacb0\uc815 \uc624\ub958",
                   "OCR \ud488\uc9c8 \ud3b8\ucc28 (\uc2a4\uce94 \ubb38\uc11c \uc624\uc778\uc2dd)"],
                  font_size=11, color=LIGHT_GRAY, bullet=True)

    # Card 3
    x3 = x2 + card_w + gap
    add_card(slide, x3, y_start, card_w, card_h)
    add_text(slide, x3 + Emu(182880), y_start + Emu(137160), card_w - Emu(365760), Emu(320040),
             "[DATA] \ub370\uc774\ud130 \ud65c\uc6a9 \ubd80\uc7ac", font_size=15, color=GOLD, bold=True)
    add_multiline(slide, x3 + Emu(182880), y_start + Emu(502920),
                  card_w - Emu(365760), card_h - Emu(640080),
                  ["\ubd84\uc11d \uacb0\uacfc\uac00 \uac1c\ubcc4 \ud30c\uc77c\ub85c \ubd84\uc0b0",
                   "\uacfc\uac70 \ubd84\uc11d \uacb0\uacfc \uac80\uc0c9 \ubd88\uac00\ub2a5",
                   "\ucd95\uc801\ub41c \uc9c0\uc2dd\uc758 \uc7ac\ud65c\uc6a9 \uc5c6\uc74c",
                   "\ud1b5\ud569 \ud1b5\uacc4/\ub300\uc2dc\ubcf4\ub4dc \ubd80\uc7ac"],
                  font_size=11, color=LIGHT_GRAY, bullet=True)

    add_text(slide, MARGIN_L, Emu(4206240), CONTENT_W, Emu(365760),
             "-> AI \uc790\ub3d9\ud654 + \uc9c0\uc2dd \ucd95\uc801 \uc2dc\uc2a4\ud15c\uc73c\ub85c \uc18d\ub3c4/\uc815\ud655\ub3c4/\uc77c\uad00\uc131/\ud65c\uc6a9\uc131 \ub3d9\uc2dc \ud574\uacb0",
             font_size=12, color=GOLD, alignment=PP_ALIGN.CENTER)


def slide_04_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "02  \uc194\ub8e8\uc158 -- 5\ub2e8\uacc4 \uc790\ub3d9 \ud30c\uc774\ud504\ub77c\uc778",
               "\uc5c5\ub85c\ub4dc \ud55c \ubc88\uc73c\ub85c OCR -> \ubd84\uc11d -> \ubcf4\uace0\uc11c -> AI \ucc57\ubd07\uae4c\uc9c0 \uc804\uc790\ub3d9")

    steps = [
        ("(1) \uc5c5\ub85c\ub4dc",
         ["PDF / IMG / HTML", "DART ZIP \uc790\ub3d9 \uac10\uc9c0"],
         ["\ub4dc\ub798\uadf8 \uc564 \ub4dc\ub86d", "\ub610\ub294 DART \ubc30\uce58"]),
        ("(2) OCR",
         ["EasyOCR \ub2e8\uc77c \uc5d4\uc9c4", "OpenCV 8\ub2e8\uacc4 \uc804\ucc98\ub9ac"],
         ["\ud55c/\uc601 \uc9c0\uc6d0", "BOM/UTF-16 \uc815\uaddc\ud654"]),
        ("(3) \uc804\ucc98\ub9ac",
         ["\uacc4\uce35\uc801 \uccad\ud0b9", "\uba54\ud0c0\ub370\uc774\ud130 \ucd94\ucd9c"],
         ["\uc7ac\ubb34 \ub370\uc774\ud130 \ubcf4\uc874", "\ub178\uc774\uc988 \ud544\ud130\ub9c1"]),
        ("(4) LLM",
         ["EXAONE 3.5 7.8B", "(\ub85c\uceec, \ube44\uc6a9 0\uc6d0)"],
         ["\uce74\ud14c\uace0\ub9ac/\uc694\uc57d/\uadfc\uac70", "\uc7ac\ubb34\uc9c0\ud45c \uc790\ub3d9 \ucd94\ucd9c"]),
        ("(5) \ucd9c\ub825",
         ["DB + PDF \ubcf4\uace0\uc11c", "BGE-M3 \uc784\ubca0\ub529"],
         ["ChromaDB \uc790\ub3d9 \uc800\uc7a5", "RAG \ucc57\ubd07 \ud65c\uc6a9"]),
    ]

    step_w = Emu(2012950)
    step_h = Emu(3200400)
    x = MARGIN_L
    y_top = Emu(1371600)
    arrow_w = Emu(228600)

    for i, (title, desc, detail) in enumerate(steps):
        add_card(slide, x, y_top, step_w, step_h)
        add_text(slide, x + Emu(91440), y_top + Emu(137160),
                 step_w - Emu(182880), Emu(320040),
                 title, font_size=14, color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER)
        add_multiline(slide, x + Emu(91440), y_top + Emu(548640),
                      step_w - Emu(182880), Emu(731520),
                      desc, font_size=10, color=WHITE,
                      alignment=PP_ALIGN.CENTER, line_spacing=1.4)
        add_divider(slide, x + Emu(274320), y_top + Emu(1371600),
                    step_w - Emu(548640), color=CARD_BORDER)
        add_multiline(slide, x + Emu(91440), y_top + Emu(1554480),
                      step_w - Emu(182880), Emu(731520),
                      detail, font_size=9, color=MID_GRAY,
                      alignment=PP_ALIGN.CENTER, line_spacing=1.4)
        x += step_w
        if i < len(steps) - 1:
            add_text(slide, x, y_top + Emu(640080), arrow_w, Emu(274320),
                     "->", font_size=18, color=GOLD,
                     alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
            x += arrow_w

    add_text(slide, MARGIN_L, Emu(4754880), CONTENT_W, Emu(365760),
             "\uc804 \uacfc\uc815 \uc790\ub3d9 -- \ubb38\uc11c \uc5c5\ub85c\ub4dc \uc989\uc2dc \ubd84\uc11d \uc2dc\uc791, \uc644\ub8cc \uc2dc \uc774\uba54\uc77c \uc54c\ub9bc \ubc1c\uc1a1",
             font_size=11, color=GOLD, alignment=PP_ALIGN.CENTER)


def slide_05_ui(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "03  \ud50c\ub7ab\ud3fc UI",
               "Monochrome Authority Theme -- \ud504\ub9ac\ubbf8\uc5c4 \uae08\uc735 \uc778\uc0ac\uc774\ud2b8 \uc778\ud130\ud398\uc774\uc2a4")

    card_w = Emu(3383280)
    card_h = Emu(1828800)
    gap = Emu(274320)
    y1 = Emu(1371600)
    y2 = y1 + card_h + gap

    panels = [
        ("\ub85c\uadf8\uc778 + \uc778\uc99d \uc2dc\uc2a4\ud15c",
         ["\u03a9 \ube0c\ub79c\ub529 + JWT \uc778\uc99d",
          "\uc774\uba54\uc77c \uc778\uc99d + \ube44\ubc00\ubc88\ud638 \ucc3e\uae30",
          "\uad00\ub9ac\uc790 \ub4f1\ub85d (Master Key)"]),
        ("\ubb38\uc11c \uc5c5\ub85c\ub4dc & \ubd84\uc11d",
         ["\ub4dc\ub798\uadf8 \uc564 \ub4dc\ub86d \uc5c5\ub85c\ub4dc",
          "\uc2e4\uc2dc\uac04 \ubd84\uc11d \uc9c4\ud589\ub960 \ud45c\uc2dc",
          "\uce74\ud14c\uace0\ub9ac\ubcc4 \ud544\ud130\ub9c1 + \uac80\uc0c9"]),
        ("\ubd84\uc11d \uacb0\uacfc \uc0c1\uc138",
         ["OCR \ud14d\uc2a4\ud2b8 + LLM \ubd84\uc11d \uacb0\uacfc",
          "Insight \uc0dd\uc131 (S/A/B/C \ub4f1\uae09)",
          "PDF \ubcf4\uace0\uc11c \ubbf8\ub9ac\ubcf4\uae30/\ub2e4\uc6b4\ub85c\ub4dc"]),
        ("\uad00\ub9ac\uc790 \ub300\uc2dc\ubcf4\ub4dc",
         ["\uc804\uccb4 \ubb38\uc11c \ud1b5\uacc4 + \ucd5c\uadfc \ud65c\ub3d9",
          "\uc720\uc800 \uad00\ub9ac (\ubaa9\ub85d / \uc5ed\ud560 \ubcc0\uacbd)",
          "\uc804\uccb4 \ubb38\uc11c \uc870\ud68c + \uc7ac\ubd84\ub958"]),
        ("AI \ucc57\ubd07 (Omega Cortex)",
         ["\uba40\ud2f0\ud134 \ub300\ud654 \uc778\ud130\ud398\uc774\uc2a4",
          "RAG \uae30\ubc18 \uadfc\uac70 \uc81c\uc2dc \ub2f5\ubcc0",
          "\uae30\uc5c5\uba85/\uc5f0\ub3c4 \uc790\ub3d9\uc644\uc131 \uac80\uc0c9"]),
        ("UI/UX \ud2b9\uc9d5",
         ["\uac80\uc815(#050505) + \uace8\ub4dc(#C0A060)",
          "Glassmorphism \uce74\ub4dc \ub514\uc790\uc778",
          "\uc5ed\ud560\ubcc4 \uba54\ub274 \ubd84\uae30 (user/admin)"]),
    ]

    for i, (title, items) in enumerate(panels):
        col = i % 3
        row = i // 3
        x = MARGIN_L + col * (card_w + gap)
        y = y1 + row * (card_h + gap)
        add_card(slide, x, y, card_w, card_h)
        add_text(slide, x + Emu(137160), y + Emu(91440),
                 card_w - Emu(274320), Emu(274320),
                 title, font_size=13, color=GOLD, bold=True)
        add_multiline(slide, x + Emu(137160), y + Emu(457200),
                      card_w - Emu(274320), card_h - Emu(548640),
                      items, font_size=10, color=LIGHT_GRAY, bullet=True)


def slide_06_differentiators(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "04  \ud575\uc2ec \ucc28\ubcc4\uc810 & \uc131\uacfc",
               "\uc65c Omega CivicFlow\uc778\uac00?")

    items = [
        ("01", "\uc804\uc790\ub3d9 \ud30c\uc774\ud504\ub77c\uc778",
         "\uc5c5\ub85c\ub4dc 1\ubc88 -> OCR -> \uc804\ucc98\ub9ac -> LLM -> PDF \ubcf4\uace0\uc11c -> \uc784\ubca0\ub529 -> \uc774\uba54\uc77c\uae4c\uc9c0 \uc804\uc790\ub3d9"),
        ("02", "Dual-LLM \uc544\ud0a4\ud14d\ucc98",
         "EXAONE 3.5(\ub85c\uceec, \ube44\uc6a9 0\uc6d0) + Gemini 2.5 Pro(\uc804\ub7b5 \uc2ec\uce35) 2\ub2e8 \uad6c\uc870"),
        ("03", "Supervisor \ub3c5\ub9bd \uac10\uc0ac",
         "Gemini 2.5 Flash\uac00 Insight \uacb0\uacfc\ub97c \ub3c5\ub9bd \uac80\uc99d -- 5\ub2e8\uacc4 \uc2e0\ub8b0\ub3c4 \uad50\uc815"),
        ("04", "300K \ubca1\ud130 RAG",
         "BGE-M3 1024-dim x 300,306 \ubca1\ud130 \uc778\ub371\uc2f1 -- \ubca1\ud130+BM25+\uba54\ud0c0 \ud558\uc774\ube0c\ub9ac\ub4dc \uac80\uc0c9"),
        ("05", "\uc804\ub7b5 \ub4f1\uae09 \uc2dc\uc2a4\ud15c",
         "Gemini 2.5 Pro -> \ud22c\uc790 \uc2dc\uc0ac\uc810/\ub9ac\uc2a4\ud06c/S/A/B/C \ub4f1\uae09 \uc790\ub3d9 \ubd80\uc5ec"),
        ("06", "3,135\uac74 \uc2e4\uc801",
         "97.2% \uc694\uc57d \uc644\uc131\ub960 (3,047/3,135) | 12,211\uac74 \uc7ac\ubb34\uc9c0\ud45c | 1,106\uac1c \uace0\uc720 \uae30\uc5c5"),
        ("07", "\uc644\uc804\ud55c \uc778\uc99d \uccb4\uacc4",
         "JWT + \uc774\uba54\uc77c \uc778\uc99d + RBAC + CORS \uc81c\ud55c + .env \ubd84\ub9ac"),
    ]

    y = Emu(1371600)
    for num, title, desc in items:
        add_text(slide, MARGIN_L, y, Emu(365760), Emu(274320),
                 num, font_size=18, color=GOLD, bold=True, font_name=FONT_EN)
        add_text(slide, Emu(1188720), y, Emu(2286000), Emu(274320),
                 title, font_size=13, color=WHITE, bold=True)
        add_text(slide, Emu(3566160), y + Emu(22860), Emu(7315200), Emu(274320),
                 desc, font_size=10, color=LIGHT_GRAY)
        if num != "07":
            add_divider(slide, MARGIN_L, y + Emu(365760),
                        CONTENT_W, color=RGBColor(0x22, 0x22, 0x22))
        y += Emu(502920)


def slide_07_backend(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "05  \ubc31\uc5d4\ub4dc \uc11c\ube44\uc2a4 \ub808\uc774\uc5b4 \uc0c1\uc138",
               "27\uac1c \uc804\ubb38 \uc11c\ube44\uc2a4 + 6\ub2e8\uacc4 LLM \uc624\ucf00\uc2a4\ud2b8\ub808\uc774\uc158")

    services = [
        ("OCR \uc5d4\uc9c4", "EasyOCR \ub2e8\uc77c \uc5d4\uc9c4 + OpenCV 8\ub2e8\uacc4 \uc774\ubbf8\uc9c0 \uc804\ucc98\ub9ac, BOM/UTF-16 \uc815\uaddc\ud654"),
        ("LLM \ubd84\uc11d", "Ollama API\ub85c EXAONE 3.5 7.8B\uc5d0 \ubb38\uc11c \ubd84\uc11d \uc694\uccad -> JSON \uad6c\uc870\ud654"),
        ("\uc804\ub7b5 Insight", "Gemini 2.5 Pro\ub85c \ud22c\uc790 \uc2dc\uc0ac\uc810/\ub9ac\uc2a4\ud06c/\uc804\ub7b5 \ub4f1\uae09 \uc0dd\uc131"),
        ("Supervisor", "Gemini 2.5 Flash\ub85c Insight \uacb0\uacfc \ub3c5\ub9bd \uac10\uc0ac/\uc2e0\ub8b0\ub3c4 \uad50\uc815"),
        ("\ud14d\uc2a4\ud2b8 \uc804\ucc98\ub9ac", "\uacc4\uce35\uc801 \ud5e4\ub354 \uc8fc\uc785 + \ubb38\uc11c \uad6c\uc870 \ucd94\ucd9c + \ub178\uc774\uc988 \ud544\ud130\ub9c1"),
        ("\ubca1\ud130 \uc784\ubca0\ub529", "BGE-M3 1024-dim + \ud558\uc774\ube0c\ub9ac\ub4dc \uac80\uc0c9 (0.40 / 0.25 / 0.35)"),
        ("PDF \ubcf4\uace0\uc11c", "\ubd84\uc11d \uacb0\uacfc\ub97c \ud55c\uae00 PDF \uc694\uc57d \ubcf4\uace0\uc11c\ub85c \uc790\ub3d9 \uc0dd\uc131"),
        ("JWT \uc778\uc99d", "HS256 \ud1a0\ud070 \ubc1c\uae09/\uac80\uc99d, bcrypt \ud574\uc2f1, \uc774\uba54\uc77c \uc778\uc99d"),
        ("AI \ucc57\ubd07", "\uba40\ud2f0\ud134 \ub300\ud654 + RAG \uae30\ubc18 \uadfc\uac70 \uc81c\uc2dc \ub2f5\ubcc0"),
        ("\uc774\uba54\uc77c \ubc1c\uc1a1", "Gmail SMTP\ub85c \ubd84\uc11d \uc644\ub8cc \uc54c\ub9bc + \uacb0\uacfc \uc694\uc57d \uc790\ub3d9 \ubc1c\uc1a1"),
    ]

    y = Emu(1280160)
    row_h = Emu(411480)
    col1_w = Emu(2103120)

    for name, desc in services:
        add_text(slide, MARGIN_L, y, col1_w, row_h,
                 name, font_size=11, color=GOLD, bold=True)
        add_text(slide, MARGIN_L + col1_w + Emu(182880), y,
                 Emu(8229600), row_h,
                 desc, font_size=10, color=LIGHT_GRAY)
        y += row_h

    add_divider(slide, MARGIN_L, y + Emu(45720), CONTENT_W, color=CARD_BORDER)
    add_text(slide, MARGIN_L, y + Emu(137160), CONTENT_W, Emu(365760),
             "6\ub2e8\uacc4 LLM \uc624\ucf00\uc2a4\ud2b8\ub808\uc774\uc158: Router -> Planner -> Judge -> Synthesizer -> Critic -> Reviser (\uc870\uac74\ubd80)",
             font_size=10, color=MID_GRAY)


def slide_08_ocr(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "06  OCR \ub0b4\ubd80 \uba54\ucee4\ub2c8\uc998",
               "EasyOCR + OpenCV 8\ub2e8\uacc4 \uc774\ubbf8\uc9c0 \uc804\ucc98\ub9ac \ud30c\uc774\ud504\ub77c\uc778")

    card_w = Emu(5120640)
    card_h = Emu(2834640)
    gap = Emu(274320)
    y_start = Emu(1371600)

    add_card(slide, MARGIN_L, y_start, card_w, card_h)
    add_text(slide, MARGIN_L + Emu(182880), y_start + Emu(137160),
             card_w - Emu(365760), Emu(320040),
             "OpenCV 8\ub2e8\uacc4 \uc804\ucc98\ub9ac \ud30c\uc774\ud504\ub77c\uc778", font_size=14, color=GOLD, bold=True)
    add_multiline(slide, MARGIN_L + Emu(182880), y_start + Emu(502920),
                  card_w - Emu(365760), card_h - Emu(640080),
                  ["1. \uadf8\ub808\uc774\uc2a4\ucf00\uc77c + RGBA \ub178\ub9d0\ub77c\uc774\uc81c\uc774\uc158",
                   "2. \ud574\uc0c1\ub3c4 \uc5c5\uc2a4\ucf00\uc77c (INTER_CUBIC, \u22652000px)",
                   "3. CLAHE \uc801\uc751\ud615 \ub300\ube44 \uac15\ud654 (clipLimit=1.5)",
                   "4. Hough \uae30\ubc18 \uc790\ub3d9 \uae30\uc6b8\uae30 \ubcf4\uc815 (Deskew \u00b110\u00b0)",
                   "5. EasyOCR \uc790\uccb4 \uc774\uc9c4\ud654 \ud65c\uc6a9 (\uc0ac\uc804 \uc774\uc9c4\ud654 X)",
                   "6. PIL Fallback (OpenCV \ubbf8\uc124\uce58 \uc2dc 3\ub2e8\uacc4 \uacbd\ub7c9 \uacbd\ub85c)",
                   "7. ZIP \ud30c\uc77c \uc790\ub3d9 \uac10\uc9c0 -> DART XBRL/XML \ud30c\uc2f1",
                   "8. DART _lab-ko.xml \ud55c\uad6d\uc5b4 \ub77c\ubca8 \uc6b0\uc120 \ud30c\uc2f1"],
                  font_size=10, color=LIGHT_GRAY, bullet=True)

    x2 = MARGIN_L + card_w + gap
    card_w2 = CONTENT_W - card_w - gap
    add_card(slide, x2, y_start, card_w2, card_h)
    add_text(slide, x2 + Emu(182880), y_start + Emu(137160),
             card_w2 - Emu(365760), Emu(320040),
             "\uc65c EasyOCR\uc778\uac00?", font_size=14, color=GOLD, bold=True)
    add_multiline(slide, x2 + Emu(182880), y_start + Emu(502920),
                  card_w2 - Emu(365760), card_h - Emu(640080),
                  ["\ud55c\uad6d\uc5b4/\uc601\uc5b4 \ub3d9\uc2dc \uc9c0\uc6d0",
                   "Pip \uc124\uce58 \ud558\ub098\ub85c \uc989\uc2dc \uad6c\ub3d9",
                   "\uacbd\ub7c9 + GPU/CPU \uc720\uc5f0 \ub300\uc751",
                   "OpenCV \uc804\ucc98\ub9ac\ub85c \ud488\uc9c8 \ubcf4\uc644",
                   "\uae08\uc735 \ubb38\uc11c \uc218\uce58/\uad04\ud638/% \ubcf4\uc874"],
                  font_size=10, color=LIGHT_GRAY, bullet=True)

    add_text(slide, MARGIN_L, y_start + card_h + Emu(182880), CONTENT_W, Emu(365760),
             "\ud488\uc9c8 \ub4f1\uae09 \uc790\ub3d9 \ubd84\ub958: OCR \uc2e0\ub8b0\ub3c4 \uae30\ubc18 good / low / very_low -> \uc800\ud488\uc9c8 \ud398\uc774\uc9c0 \uc790\ub3d9 \uc81c\uc678",
             font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def slide_09_llm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "07  LLM \ubd84\uc11d \uc5d4\uc9c4 \ub0b4\ubd80 \uad6c\uc870",
               "EXAONE 3.5 7.8B (LG AI Research) -- Ollama \ub85c\uceec \uc2e4\ud589")

    card_w = Emu(3383280)
    card_h = Emu(3200400)
    gap = Emu(274320)
    y_start = Emu(1371600)

    # Card 1
    add_card(slide, MARGIN_L, y_start, card_w, card_h)
    add_text(slide, MARGIN_L + Emu(182880), y_start + Emu(137160),
             card_w - Emu(365760), Emu(320040),
             "\ubaa8\ub378 \uad6c\uc131 & \ucd9c\ub825", font_size=13, color=GOLD, bold=True)
    add_multiline(slide, MARGIN_L + Emu(182880), y_start + Emu(502920),
                  card_w - Emu(365760), card_h - Emu(640080),
                  ["\ubca0\uc774\uc2a4: EXAONE 3.5 7.8B (LG AI)",
                   "Ollama \ub85c\uceec (\ube44\uc6a9 0\uc6d0, \uac1c\uc778\uc815\ubcf4 \ubcf4\ud638)",
                   "",
                   "\ubd84\uc11d \ucd9c\ub825 (JSON \uad6c\uc870\ud654):",
                   "  summary -- \ubb38\uc11c \uc694\uc57d",
                   "  category -- \ubb38\uc11c \uc720\ud615 \ubd84\ub958",
                   "  company_name -- \uae30\uc5c5\uba85 \ucd94\ucd9c",
                   "  financial_metrics -- \uc7ac\ubb34 \uc9c0\ud45c",
                   "  insight_vectors -- \ud575\uc2ec \uc2dc\uc0ac\uc810",
                   "  evidence -- \ubd84\uc11d \uadfc\uac70 \uc6d0\ubb38",
                   "  key_points -- \ud575\uc2ec \ud3ec\uc778\ud2b8"],
                  font_size=9, color=LIGHT_GRAY, bullet=True)

    # Card 2
    x2 = MARGIN_L + card_w + gap
    add_card(slide, x2, y_start, card_w, Emu(1554480))
    add_text(slide, x2 + Emu(182880), y_start + Emu(137160),
             card_w - Emu(365760), Emu(320040),
             "\ubaa8\ub378 \ud30c\ub77c\ubbf8\ud130", font_size=13, color=GOLD, bold=True)
    add_multiline(slide, x2 + Emu(182880), y_start + Emu(502920),
                  card_w - Emu(365760), Emu(914400),
                  ["temperature: 0.1 (\uacb0\uc815\uc801, \ud658\uac01 \ucd5c\uc18c)",
                   "seed: 42 (\uc7ac\ud604\uc131 \ud655\ubcf4)",
                   "top_p: 0.9 (\ub2e4\uc591\uc131/\uc815\ud655\ub3c4 \uade0\ud615)",
                   "num_ctx: 16384 (\uae34 \ubb38\uc11c)",
                   "num_predict: 3072~4096 (\uc644\uc131\ub3c4 \ud655\ubcf4)",
                   "repeat_penalty: 1.05 (\uc870\uae30 \uc885\ub8cc \ubc29\uc9c0)"],
                  font_size=10, color=LIGHT_GRAY, bullet=True)

    # Card 3
    x3 = x2 + card_w + gap
    add_card(slide, x3, y_start, card_w, card_h)
    add_text(slide, x3 + Emu(182880), y_start + Emu(137160),
             card_w - Emu(365760), Emu(320040),
             "\uc65c EXAONE 3.5\uc778\uac00?", font_size=13, color=GOLD, bold=True)
    add_multiline(slide, x3 + Emu(182880), y_start + Emu(502920),
                  card_w - Emu(365760), card_h - Emu(640080),
                  ["LG AI Research \ud55c\uad6d\uc5b4 \ud2b9\ud654",
                   "\ud55c\uad6d\uc5b4 \uc131\ub2a5 7B\uae09 \uc0c1\uc704\uad8c",
                   "\uc911\uad6d\uc5b4 \ud63c\uc785 \ubb38\uc81c \uc5c6\uc74c (!=Qwen)",
                   "Ollama \ub85c\uceec -> API \ube44\uc6a9 0\uc6d0",
                   "\uac1c\uc778\uc815\ubcf4 \ubcf4\ud638 (\uc678\ubd80 \uc804\uc1a1 \uc5c6\uc74c)",
                   "\ucee4\uc2a4\ud140 \uc2dc\uc2a4\ud15c \ud504\ub86c\ud504\ud2b8 \uc801\uc6a9",
                   "",
                   "-> \uc7ac\ubb34 \ubd84\uc11d \uc678 \ubaa8\ub4e0 \uc791\uc5c5 \ub2f4\ub2f9",
                   "  (OCR \ubcf4\uc815, RAG, \ucc57\ubd07, \uc694\uc57d)"],
                  font_size=9, color=LIGHT_GRAY, bullet=True)


def slide_10_insight(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "08  Insight \uc5d4\uc9c4 + Omega-Prime Supervisor",
               "Gemini 2.5 Pro (\uc804\ub7b5 \ubd84\uc11d) + Gemini 2.5 Flash (\ub3c5\ub9bd \uac10\uc0ac)")

    card_w = Emu(5120640)
    gap = Emu(274320)
    y_start = Emu(1371600)

    add_card(slide, MARGIN_L, y_start, card_w, Emu(3108960))
    add_text(slide, MARGIN_L + Emu(182880), y_start + Emu(137160),
             card_w - Emu(365760), Emu(320040),
             "Gemini 2.5 Pro -- Insight \uc5d4\uc9c4", font_size=14, color=GOLD, bold=True)
    add_multiline(slide, MARGIN_L + Emu(182880), y_start + Emu(502920),
                  card_w - Emu(365760), Emu(2468880),
                  ["1. Investment Thesis -- \ud575\uc2ec \ud22c\uc790 \ud310\ub2e8 \uadfc\uac70",
                   "2. Market Context -- \uc0b0\uc5c5 \ub3d9\ud5a5, \uacbd\uc7c1\uc0ac \ube44\uad50",
                   "3. Risk Factors -- \uc7ac\ubb34 \uc704\ud5d8, \uaddc\uc81c \ub9ac\uc2a4\ud06c",
                   "4. Strategic Action -- \uad6c\uccb4\uc801 \ud589\ub3d9 \uc9c0\uce68",
                   "5. Strategy Rating -- S(\uac15\ub825)/A(\uae0d\uc815)/B(\uc911\ub9bd)/C(\uc8fc\uc758)",
                   "",
                   "\uc120\ud0dd \uadfc\uac70: 100\ub9cc \ud1a0\ud070, \ud55c\uad6d\uc5b4 JSON \uc900\uc218\ub960 \uc6b0\uc218",
                   "\ubcc4\ub3c4 GCP \ud504\ub85c\uc81d\ud2b8 (4\uac1c \ud0a4 \ud480 + 429 backoff)"],
                  font_size=10, color=LIGHT_GRAY, bullet=True)

    x2 = MARGIN_L + card_w + gap
    card_w2 = CONTENT_W - card_w - gap
    add_card(slide, x2, y_start, card_w2, Emu(3108960))
    add_text(slide, x2 + Emu(182880), y_start + Emu(137160),
             card_w2 - Emu(365760), Emu(320040),
             "Gemini 2.5 Flash -- Supervisor", font_size=14, color=GOLD, bold=True)
    add_multiline(slide, x2 + Emu(182880), y_start + Emu(502920),
                  card_w2 - Emu(365760), Emu(2468880),
                  ["Omega-Prime 5\ub2e8\uacc4 \ud504\ub85c\ud1a0\ucf5c:",
                   "",
                   "1. DECOMPOSE",
                   "   \ubcc0\uc218/\uc81c\uc57d/\ubbf8\uc9c0\uc218 \ubd84\ub9ac",
                   "2. CAUSAL CHECK",
                   "   \uc778\uacfc \ubc29\ud5a5 \uac80\uc99d",
                   "3. HIDDEN RISK SCAN",
                   "   \uc228\uaca8\uc9c4 \uc704\ud5d8 \ud0d0\uc0c9",
                   "4. COUNTERFACTUAL",
                   "   \ubc18\uc0ac\uc2e4 \uc2a4\ud2b8\ub808\uc2a4 \ud14c\uc2a4\ud2b8",
                   "5. CONFIDENCE CALIBRATION",
                   "   \uc2e0\ub8b0\ub3c4 \ub808\ubca8 \ubd80\uc5ec",
                   "",
                   "\ubcc4\ub3c4 GCP \ud504\ub85c\uc81d\ud2b8 \ubd84\ub9ac \uc6b4\uc601"],
                  font_size=9, color=LIGHT_GRAY, bullet=True)

    add_text(slide, MARGIN_L, Emu(4663440), CONTENT_W, Emu(365760),
             "\uc2e0\ub8b0\ub3c4: AXIOM [99%] | CONSENSUS [85-95%] | INFERENCE [65-84%] | SPECULATION [40-64%] | EXPLORATION [<40%]",
             font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER, font_name=FONT_EN)


def slide_11_rag(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "09  RAG \uac80\uc0c9 \uc2dc\uc2a4\ud15c",
               "BGE-M3 1024-dim + \ud558\uc774\ube0c\ub9ac\ub4dc \uac80\uc0c9 + CrossEncoder Reranking")

    card_w = Emu(3383280)
    card_h = Emu(2834640)
    gap = Emu(274320)
    y_start = Emu(1371600)

    add_card(slide, MARGIN_L, y_start, card_w, card_h)
    add_text(slide, MARGIN_L + Emu(182880), y_start + Emu(137160),
             card_w - Emu(365760), Emu(320040),
             "\uc784\ubca0\ub529 \ud30c\uc774\ud504\ub77c\uc778", font_size=13, color=GOLD, bold=True)
    add_multiline(slide, MARGIN_L + Emu(182880), y_start + Emu(502920),
                  card_w - Emu(365760), card_h - Emu(640080),
                  ["\ubaa8\ub378: BAAI/bge-m3 (1024-dim)",
                   "ChromaDB: 300,306 \ubca1\ud130 \uc778\ub371\uc2f1 \uc644\ub8cc",
                   "(SQL document_chunks: 284,146 rows)",
                   "Chunk: 1,000\uc790 / Overlap: 150\uc790",
                   "A100 40GB \ud074\ub77c\uc6b0\ub4dc (10-15\ubd84)",
                   "",
                   "\uc120\ud0dd \uadfc\uac70:",
                   "MIRACL/MTEB \ud55c\uad6d\uc5b4 \ubca4\uce58\ub9c8\ud06c",
                   "OpenAI text-embedding-3 \ub300\ube44 \uc6b0\uc218"],
                  font_size=10, color=LIGHT_GRAY, bullet=True)

    x2 = MARGIN_L + card_w + gap
    add_card(slide, x2, y_start, card_w, card_h)
    add_text(slide, x2 + Emu(182880), y_start + Emu(137160),
             card_w - Emu(365760), Emu(320040),
             "\ud558\uc774\ube0c\ub9ac\ub4dc \uac80\uc0c9 \ud750\ub984", font_size=13, color=GOLD, bold=True)
    add_multiline(slide, x2 + Emu(182880), y_start + Emu(502920),
                  card_w - Emu(365760), card_h - Emu(640080),
                  ["\ubca1\ud130 \uc720\uc0ac\ub3c4: 40%",
                   "BM25 \ud0a4\uc6cc\ub4dc: 25%",
                   "\uba54\ud0c0\ub370\uc774\ud130 \ud544\ud130: 35%",
                   "  (\ud68c\uc0ac\uba85/\uc5f0\ub3c4/\uc139\uc158 \uc720\ud615)",
                   "",
                   "\uc5d4\ud2b8\ub85c\ud53c \ud544\ud130: 0.10 (noise \ucc28\ub2e8)",
                   "\uc911\ubcf5 \uc81c\uac70: 0.80 threshold",
                   "CrossEncoder Reranking \uc801\uc6a9"],
                  font_size=10, color=LIGHT_GRAY, bullet=True)

    x3 = x2 + card_w + gap
    add_card(slide, x3, y_start, card_w, card_h)
    add_text(slide, x3 + Emu(182880), y_start + Emu(137160),
             card_w - Emu(365760), Emu(320040),
             "\uc65c RAG\uc778\uac00?", font_size=13, color=GOLD, bold=True)
    add_multiline(slide, x3 + Emu(182880), y_start + Emu(502920),
                  card_w - Emu(365760), card_h - Emu(640080),
                  ["3,135\uac74 \ubd84\uc11d \ubb38\uc11c \ucc38\uc870 \ud544\uc694",
                   "LLM \ub2e8\ub3c5\uc73c\ub85c\ub294 \ubb38\uc11c DB \uac80\uc0c9 \ubd88\uac00",
                   "",
                   "\ud658\uac01 \ubc29\uc9c0:",
                   "\uadfc\uac70 \uccad\ud06c \uc9c1\uc811 \uc81c\uc2dc",
                   "Judge/Critic \ub2e8\uacc4 \uac80\uc99d",
                   "",
                   "\uc790\ub3d9 \ud655\uc7a5:",
                   "\ubb38\uc11c \ucd94\uac00 \uc2dc \uc790\ub3d9 \uc784\ubca0\ub529",
                   "\uad6c\uc870\ud654 \uccad\ud06c \uba54\ud0c0 \ud544\ud130\ub9c1"],
                  font_size=10, color=LIGHT_GRAY, bullet=True)


def slide_12_frontend(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "10  \ud504\ub860\ud2b8\uc5d4\ub4dc \uad6c\ud604 \uc0c1\uc138",
               "React 18 + Vite 6 -- 14 Pages, 4 Global Components")

    card_w = Emu(2560320)
    card_h = Emu(2194560)
    gap = Emu(182880)
    y_start = Emu(1371600)

    groups = [
        ("Public (8 Pages)",
         ["LoginPage -- \ub85c\uadf8\uc778",
          "RegisterPage -- \ud68c\uc6d0\uac00\uc785",
          "VerifyEmail -- \uc774\uba54\uc77c \uc778\uc99d",
          "ForgotPassword",
          "ResetPassword",
          "AdminRegisterPage",
          "VerifyPasswordChange",
          "HomePage -- \ub79c\ub529"]),
        ("Protected (3 Pages)",
         ["UploadPage",
          "  \ub4dc\ub798\uadf8 \uc564 \ub4dc\ub86d \uc5c5\ub85c\ub4dc",
          "  \uc2e4\uc2dc\uac04 \ubd84\uc11d \uc9c4\ud589\ub960",
          "MyPage",
          "  \ub0b4 \ubb38\uc11c + \uce74\ud14c\uace0\ub9ac \ud1b5\uacc4",
          "DocumentDetail",
          "  \ubd84\uc11d \uacb0\uacfc + Insight",
          "  PDF \ubcf4\uace0\uc11c \ubbf8\ub9ac\ubcf4\uae30"]),
        ("Admin Only (3 Pages)",
         ["AdminDashboard",
          "  \uc804\uccb4 \ubb38\uc11c \ud1b5\uacc4/\ud65c\ub3d9",
          "AdminUsers",
          "  \uc720\uc800 \ubaa9\ub85d/\uc5ed\ud560 \ubcc0\uacbd",
          "AdminDocuments",
          "  \uc804\uccb4 \ubb38\uc11c \uc870\ud68c/\uc7ac\ubd84\ub958"]),
        ("\uc804\uc5ed \ucef4\ud3ec\ub10c\ud2b8",
         ["Navbar (\uc5ed\ud560\ubcc4 \uba54\ub274)",
          "ChatBot (AI \ucc57\ubd07 \ud50c\ub85c\ud305)",
          "ProtectedRoute (JWT \uac80\uc99d)",
          "SideDecorations",
          "",
          "AuthContext (\uc804\uc5ed \uc778\uc99d)",
          "React Router 6 (SPA)",
          "Axios HTTP \ud074\ub77c\uc774\uc5b8\ud2b8"]),
    ]

    for i, (title, items) in enumerate(groups):
        x = MARGIN_L + i * (card_w + gap)
        add_card(slide, x, y_start, card_w, card_h)
        add_text(slide, x + Emu(137160), y_start + Emu(91440),
                 card_w - Emu(274320), Emu(274320),
                 title, font_size=12, color=GOLD, bold=True)
        add_multiline(slide, x + Emu(137160), y_start + Emu(411480),
                      card_w - Emu(274320), card_h - Emu(502920),
                      items, font_size=9, color=LIGHT_GRAY, bullet=True)

    add_text(slide, MARGIN_L, Emu(3749040), CONTENT_W, Emu(365760),
             "\ub514\uc790\uc778: Monochrome Authority -- \uac80\uc815(#050505) + \uace8\ub4dc(#C0A060), Glassmorphism \uce74\ub4dc",
             font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def slide_13_database(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "11  \ub370\uc774\ud130\ubca0\uc774\uc2a4 \uc2a4\ud0a4\ub9c8",
               "SQLite + SQLAlchemy ORM -- 11 Tables")

    tables = [
        ("users", "id, email(UK), username, password_hash, role(user|admin), is_verified",
         "\uc778\uc99d/\uad8c\ud55c \uae30\uc900"),
        ("documents", "id, user_id(FK), filename, status, report_path",
         "\ud30c\uc774\ud504\ub77c\uc778 \uc0c1\ud0dc"),
        ("pages", "id, document_id(FK), page_number, image_path",
         "PDF \ud398\uc774\uc9c0 \uc774\ubbf8\uc9c0"),
        ("ocr_texts", "id, document_id(FK), page_id(FK), raw_text, cleaned_text, confidence",
         "\uc6d0\ubcf8/\uc815\uc81c + \uc2e0\ub8b0\ub3c4"),
        ("analysis_results", "id, document_id(FK), summary, category, financial_metrics, evidence",
         "EXAONE \ubd84\uc11d \uacb0\uacfc"),
        ("document_insights", "id, document_id(FK), investment_thesis, risk_factors, strategy_rating",
         "Gemini Insight+Supervisor"),
        ("document_metadata", "id, document_id(FK), company_name_norm, corp_code, fiscal_year",
         "DART \uacf5\uc2dc \uba54\ud0c0"),
        ("document_chunks", "chunk_uid, document_id(FK), text, token_count, vector_collection",
         "RAG \uac80\uc0c9 \ub2e8\uc704"),
        ("financial_facts", "fact_uid, document_id(FK), metric_name, metric_value_num",
         "\uad6c\uc870\ud654 \uc7ac\ubb34 \uc9c0\ud45c"),
        ("company_profiles", "company_name_norm(UK), corp_code, latest_fiscal_year",
         "\uae30\uc5c5 \ud504\ub85c\ud544 \uce90\uc2dc"),
        ("reclassifications", "id, document_id(FK), reclassified_by(FK), reason",
         "\uc7ac\ubd84\ub958 \uac10\uc0ac"),
    ]

    y = Emu(1234440)
    row_h = Emu(365760)

    add_text(slide, MARGIN_L, y, Emu(1828800), Emu(274320),
             "Table", font_size=11, color=GOLD, bold=True, font_name=FONT_EN)
    add_text(slide, MARGIN_L + Emu(1920240), y, Emu(6217920), Emu(274320),
             "Primary Columns", font_size=11, color=GOLD, bold=True, font_name=FONT_EN)
    add_text(slide, MARGIN_L + Emu(8229600), y, Emu(2560320), Emu(274320),
             "\uc5ed\ud560", font_size=11, color=GOLD, bold=True)
    y += Emu(320040)
    add_divider(slide, MARGIN_L, y, CONTENT_W, color=CARD_BORDER)
    y += Emu(91440)

    for name, cols, role in tables:
        add_text(slide, MARGIN_L, y, Emu(1828800), row_h,
                 name, font_size=9, color=WHITE, bold=True, font_name=FONT_EN)
        add_text(slide, MARGIN_L + Emu(1920240), y, Emu(6217920), row_h,
                 cols, font_size=8, color=LIGHT_GRAY, font_name=FONT_EN)
        add_text(slide, MARGIN_L + Emu(8229600), y, Emu(2560320), row_h,
                 role, font_size=9, color=MID_GRAY)
        y += row_h

    add_text(slide, MARGIN_L, y + Emu(91440), CONTENT_W, Emu(365760),
             "User 1:N -> Document 1:N -> Page / OcrText / AnalysisResult / Insight / Metadata / Chunk / Fact",
             font_size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER, font_name=FONT_EN)


def slide_14_security(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "12  \ubcf4\uc548 & \uc778\uc99d \uccb4\uacc4",
               "JWT + \uc774\uba54\uc77c \uc778\uc99d + \uc5ed\ud560 \uae30\ubc18 \uc811\uadfc \uc81c\uc5b4 + \ub370\uc774\ud130 \ubcf4\ud638")

    card_w = Emu(2560320)
    card_h = Emu(2834640)
    gap = Emu(182880)
    y_start = Emu(1371600)

    cards = [
        ("\uc778\uc99d \uba54\ucee4\ub2c8\uc998",
         ["JWT (HS256) \ud1a0\ud070 \uae30\ubc18",
          "  \uc561\uc138\uc2a4 \ud1a0\ud070 24\uc2dc\uac04 \ub9cc\ub8cc",
          "bcrypt \ud574\uc2f1 (salt \uc790\ub3d9)",
          "\uc774\uba54\uc77c \uc778\uc99d \ud544\uc218 (SMTP)",
          "\ube44\ubc00\ubc88\ud638 \ubcc0\uacbd \uc2dc \uc774\uba54\uc77c \uc7ac\ud655\uc778"]),
        ("\uc5ed\ud560 \uae30\ubc18 \uc811\uadfc \uc81c\uc5b4 (RBAC)",
         ["user -- \uc77c\ubc18 \uc0ac\uc6a9\uc790",
          "  \ubb38\uc11c \uc5c5\ub85c\ub4dc/\uc870\ud68c/\ubd84\uc11d",
          "  \ubcf8\uc778 \ubb38\uc11c\ub9cc \uc811\uadfc \uac00\ub2a5",
          "admin -- \uad00\ub9ac\uc790",
          "  \uc804\uccb4 \ubb38\uc11c/\uc720\uc800 \uad00\ub9ac",
          "  \uce74\ud14c\uace0\ub9ac \uc7ac\ubd84\ub958/\ud1b5\uacc4"]),
        ("\ub370\uc774\ud130 \ubcf4\ud638",
         ["GCP \uc11c\ube44\uc2a4 \ud0a4 -> .gitignore",
          ".env \ud658\uacbd\ubcc0\uc218 \ubd84\ub9ac",
          "CORS \ucd9c\ucc98 \uc81c\ud55c",
          "\ud30c\uc77c \ud06c\uae30 \uc81c\ud55c (700MB \uc0c1\ud55c)",
          "DB \uc678\ubd80 \ubd84\ub9ac \uc800\uc7a5"]),
        ("\uc65c JWT\uc778\uac00?",
         ["Stateless -- \uc11c\ubc84 \uc138\uc158 \uc5c6\uc74c",
          "FastAPI \ube44\ub3d9\uae30\uc5d0 \ucd5c\uc801",
          "DB \uc138\uc158 \ud14c\uc774\ube14 \ubd88\ud544\uc694",
          "SPA \ub77c\uc6b0\ud305 \uac04 \uc778\uc99d \uc720\uc9c0",
          "24\uc2dc\uac04 \ub9cc\ub8cc -> \ud0c8\ucde8 \uc704\ud5d8 \ucd5c\uc18c"]),
    ]

    for i, (title, items) in enumerate(cards):
        x = MARGIN_L + i * (card_w + gap)
        add_card(slide, x, y_start, card_w, card_h)
        add_text(slide, x + Emu(137160), y_start + Emu(91440),
                 card_w - Emu(274320), Emu(274320),
                 title, font_size=12, color=GOLD, bold=True)
        add_multiline(slide, x + Emu(137160), y_start + Emu(411480),
                      card_w - Emu(274320), card_h - Emu(502920),
                      items, font_size=9, color=LIGHT_GRAY, bullet=True)


def slide_15_techstack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "13  \uae30\uc220 \uc2a4\ud0dd \uc694\uc57d",
               "Technology Stack Overview")

    stack = [
        ("Frontend", "React 18 + Vite 6 + React Router 6", "SPA, Lucide Icons, Axios"),
        ("Backend", "FastAPI 0.135 + Uvicorn 0.41", "\ube44\ub3d9\uae30 ASGI, Swagger \uc790\ub3d9 \ubb38\uc11c"),
        ("Database", "SQLite + SQLAlchemy 2.0 + Alembic", "\ub2e8\uc77c \ud30c\uc77c DB, ORM, \ub9c8\uc774\uadf8\ub808\uc774\uc158"),
        ("Vector DB", "ChromaDB + BGE-M3 (1024-dim)", "MIRACL/MTEB \ud55c\uad6d\uc5b4 \uc0c1\uc704"),
        ("\ub85c\uceec LLM", "Ollama -- EXAONE 3.5 7.8B", "\ud55c\uad6d\uc5b4 \ud2b9\ud654, \ube44\uc6a9 0\uc6d0"),
        ("\ud074\ub77c\uc6b0\ub4dc AI", "Gemini 2.5 Pro + Flash (GCP)", "\ubcc4\ub3c4 GCP \ud504\ub85c\uc81d\ud2b8 \ub3c5\ub9bd \uc6b4\uc601"),
        ("OCR", "EasyOCR + OpenCV 8\ub2e8\uacc4 \uc804\ucc98\ub9ac", "\ud55c/\uc601 \uc9c0\uc6d0, Hough Deskew, CLAHE"),
        ("Reranker", "bge-reranker-v2-m3-ko (CrossEncoder)", "\uac80\uc0c9 \uacb0\uacfc \uc815\ubc00\ub3c4 \ud5a5\uc0c1"),
        ("\ube44\ub3d9\uae30", "Celery + Redis", "\ubc31\uadf8\ub77c\uc6b4\ub4dc OCR/\ubd84\uc11d/PDF"),
        ("\uc778\uc99d", "JWT (HS256) + bcrypt", "Stateless, SPA \ucd5c\uc801"),
        ("\uc774\uba54\uc77c", "Gmail SMTP", "\ubd84\uc11d \uc644\ub8cc \uc790\ub3d9 \uc54c\ub9bc"),
        ("\uc678\ubd80 API", "DART OpenAPI", "\uae08\uac10\uc6d0 \uacf5\uc2dc \ub370\uc774\ud130 \uc870\ud68c"),
        ("GPU Cloud", "RunPod / Lambda / Vast.ai (A100)", "\uc784\ubca0\ub529 \uc804\uc6a9 (sm_120 \uc6b0\ud68c)"),
    ]

    y = Emu(1234440)
    row_h = Emu(365760)

    add_text(slide, MARGIN_L, y, Emu(1554480), Emu(274320),
             "\uad6c\ubd84", font_size=11, color=GOLD, bold=True)
    add_text(slide, MARGIN_L + Emu(1645920), y, Emu(4572000), Emu(274320),
             "\uae30\uc220", font_size=11, color=GOLD, bold=True)
    add_text(slide, MARGIN_L + Emu(6309360), y, Emu(4389120), Emu(274320),
             "\uc120\ud0dd \uadfc\uac70", font_size=11, color=GOLD, bold=True)
    y += Emu(320040)
    add_divider(slide, MARGIN_L, y, CONTENT_W, color=CARD_BORDER)
    y += Emu(68580)

    for category, tech, reason in stack:
        add_text(slide, MARGIN_L, y, Emu(1554480), row_h,
                 category, font_size=9, color=WHITE, bold=True)
        add_text(slide, MARGIN_L + Emu(1645920), y, Emu(4572000), row_h,
                 tech, font_size=9, color=LIGHT_GRAY, font_name=FONT_EN)
        add_text(slide, MARGIN_L + Emu(6309360), y, Emu(4389120), row_h,
                 reason, font_size=9, color=MID_GRAY)
        y += row_h


def slide_16_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "13  \uc2dc\uc2a4\ud15c \uc544\ud0a4\ud14d\ucc98",
               "3-Tier Full-Stack + Dual-LLM + Supervisor")

    card_w = Emu(2468880)
    card_h = Emu(2651760)
    gap = Emu(182880)
    y_start = Emu(1371600)

    tiers = [
        ("Frontend -- React 18",
         ["Vite 6 \ube4c\ub4dc",
          "React Router 6 (SPA)",
          "AuthContext (\uc804\uc5ed \uc778\uc99d)",
          "14 Pages + 4 Components",
          "Monochrome Authority \ud14c\ub9c8"]),
        ("Backend -- FastAPI",
         ["Uvicorn ASGI \uc11c\ubc84",
          "4 Routers (RESTful API)",
          "27 Services (SRP)",
          "6\ub2e8\uacc4 LLM \uc624\ucf00\uc2a4\ud2b8\ub808\uc774\uc158",
          "Celery \ube44\ub3d9\uae30 \uc791\uc5c5"]),
        ("Data Layer",
         ["SQLite + SQLAlchemy ORM",
          "  (11 Tables, 3,135 docs)",
          "ChromaDB \ubca1\ud130 DB",
          "  (300,306 \ubca1\ud130, 1024-dim)",
          "File System (uploads/)"]),
        ("LLM Layer",
         ["EXAONE 3.5 (Ollama \ub85c\uceec)",
          "  -> OCR/RAG/\ucc57\ubd07/\uc694\uc57d",
          "Gemini 2.5 Pro (GCP)",
          "  -> \uc804\ub7b5 Insight \uc804\uc6a9",
          "Gemini 2.5 Flash (GCP)",
          "  -> Supervisor \ub3c5\ub9bd \uac10\uc0ac"]),
    ]

    for i, (title, items) in enumerate(tiers):
        x = MARGIN_L + i * (card_w + gap)
        add_card(slide, x, y_start, card_w, card_h)
        add_text(slide, x + Emu(137160), y_start + Emu(91440),
                 card_w - Emu(274320), Emu(274320),
                 title, font_size=11, color=GOLD, bold=True)
        add_multiline(slide, x + Emu(137160), y_start + Emu(411480),
                      card_w - Emu(274320), card_h - Emu(502920),
                      items, font_size=9, color=LIGHT_GRAY, bullet=True)

    add_text(slide, MARGIN_L, Emu(4206240), CONTENT_W, Emu(274320),
             "\uc678\ubd80 \uc11c\ube44\uc2a4 \uc5f0\ub3d9", font_size=12, color=GOLD, bold=True)
    add_text(slide, MARGIN_L, Emu(4480560), CONTENT_W, Emu(365760),
             "DART OpenAPI (\uae08\uac10\uc6d0)  |  Gmail SMTP (\uc54c\ub9bc)  |  A100 GPU Cloud (\uc784\ubca0\ub529)  |  Redis (\uc138\uc158 \uce90\uc2dc)",
             font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def slide_17_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, Emu(0), Emu(1645920), SLIDE_W, Emu(822960),
             "\u03A9", font_size=72, color=GOLD, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, Emu(0), Emu(2743200), SLIDE_W, Emu(640080),
             "OMEGA CIVICFLOW v4", font_size=36, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_divider(slide, Emu(4114800), Emu(3566160), Emu(3962095))
    add_text(slide, Emu(0), Emu(3749040), SLIDE_W, Emu(457200),
             "DART \uae08\uc735 \uacf5\uc2dc \ubb38\uc11c\uc758 \uc5d4\ud2b8\ub85c\ud53c\ub97c \uc18c\uac01\ud558\uace0",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text(slide, Emu(0), Emu(4206240), SLIDE_W, Emu(457200),
             "\uc804\ub7b5\uc801 \uc778\uc0ac\uc774\ud2b8\ub97c \uc0dd\uc131\ud558\ub294 \uc9c0\ub2a5\ud615 \ud50c\ub7ab\ud3fc",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text(slide, Emu(0), Emu(5029200), SLIDE_W, Emu(457200),
             "\uac10\uc0ac\ud569\ub2c8\ub2e4  |  Q & A",
             font_size=14, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_toc(prs)
    slide_03_problem(prs)
    slide_04_pipeline(prs)
    slide_05_ui(prs)
    slide_06_differentiators(prs)
    slide_07_backend(prs)
    slide_08_ocr(prs)
    slide_09_llm(prs)
    slide_10_insight(prs)
    slide_11_rag(prs)
    slide_12_frontend(prs)
    slide_13_database(prs)
    slide_14_security(prs)
    slide_15_techstack(prs)
    slide_16_architecture(prs)
    slide_17_closing(prs)

    out_path = os.path.join(
        os.path.dirname(os.path.abspath(__file__)), "..",
        "Omega_CivicFlow_v4_accurate.pptx"
    )
    out_path = os.path.normpath(out_path)
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
