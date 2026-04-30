"""PPT content extractor - file output"""
from pptx import Presentation
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

prs = Presentation(r'c:\Users\hibou\Downloads\Omega_CivicFlow_v4_updated.pptx')

output = []
output.append(f"Total slides: {len(prs.slides)}\n")

for i, slide in enumerate(prs.slides, 1):
    layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
    output.append(f"{'='*60}")
    output.append(f"SLIDE {i} (Layout: {layout_name})")
    output.append(f"{'='*60}")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    bold = False
                    font_size = None
                    for run in para.runs:
                        if run.font.size:
                            font_size = run.font.size.pt
                        if run.font.bold:
                            bold = True
                    prefix = "[B] " if bold else "    "
                    size_str = f" ({font_size}pt)" if font_size else ""
                    output.append(f"{prefix}{text}{size_str}")
        if shape.has_table:
            table = shape.table
            output.append(f"  [TABLE {len(list(table.rows))} x {len(table.columns)}]")
            for row in table.rows:
                cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                output.append(f"    | {' | '.join(cells)} |")
    output.append("")

with open(r'c:\Users\hibou\Omega_CivicFlow_v4\tools\pptx_content.txt', 'w', encoding='utf-8') as f:
    f.write('\n'.join(output))

print("Done - saved to pptx_content.txt")
