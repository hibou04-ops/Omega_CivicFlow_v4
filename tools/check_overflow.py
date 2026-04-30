"""텍스트 오버플로우 체크 — 텍스트가 박스 밖으로 튀어나오는지 확인"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Emu

prs = Presentation(r'C:\Users\hibou\Downloads\Omega_CivicFlow_v4_updated.pptx')

issues = []
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            # 텍스트 길이 체크
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                # 엄청 긴 줄 감지
                if len(txt) > 100:
                    # 수정된 텍스트인지 확인
                    if any(k in txt for k in ['32B', '2,518', 'vLLM', '하이브리드', '21개']):
                        issues.append(f"Slide {i+1}: [{len(txt)}자] {txt[:120]}...")

if issues:
    print(f"⚠ 긴 텍스트 {len(issues)}건 (오버플로 가능):")
    for iss in issues:
        print(f"  {iss}")
else:
    print("✅ 오버플로 없음")

# 원본 vs 수정 텍스트 길이 비교
from pptx import Presentation as P2
orig = P2(r'C:\Users\hibou\Downloads\Omega_CivicFlow_v4.pptx')

print("\n=== 변경된 텍스트 길이 비교 ===")
for si in range(len(prs.slides)):
    slide_new = prs.slides[si]
    slide_old = orig.slides[si]
    for sj, (sn, so) in enumerate(zip(slide_new.shapes, slide_old.shapes)):
        if sn.has_text_frame and so.has_text_frame:
            nt = sn.text_frame.text
            ot = so.text_frame.text
            if nt != ot:
                diff = len(nt) - len(ot)
                if abs(diff) > 5:
                    sign = "+" if diff > 0 else ""
                    print(f"  Slide {si+1}: {sign}{diff}자 | {ot[:60].strip()} → {nt[:60].strip()}")
