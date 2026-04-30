"""Verify PPT changes"""
from pptx import Presentation

SRC = r'c:\Users\hibou\Downloads\Omega_CivicFlow_v4_final.pptx'
prs = Presentation(SRC)
keywords = ['PostgreSQL', 'Supervisor', 'Flash', '24개', '2단']

found = []
for i, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                for kw in keywords:
                    if kw in t:
                        found.append(f"Slide {i}: [{kw}] {t[:120]}")
                        break

with open(r'c:\Users\hibou\Omega_CivicFlow_v4\tools\verify_log.txt', 'w', encoding='utf-8') as f:
    for line in found:
        f.write(line + '\n')
    f.write(f'\nTotal matches: {len(found)}\n')

print(f"Done - {len(found)} changes verified")
