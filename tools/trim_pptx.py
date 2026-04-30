"""PPTX 최종 정리 — 길어진 텍스트 간결화"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

DST = r'C:\Users\hibou\Downloads\Omega_CivicFlow_v4_updated.pptx'
prs = Presentation(DST)

TRIM = [
    # Slide 9: 모델 설명 간결화
    ("베이스: Qwen 2.5 32B-AWQ (GPU) + 7B Q4_K_M (CPU)",
     "베이스: Qwen 2.5 32B-AWQ + 7B Q4_K_M"),
    ("양자화: AWQ 4-bit (32B, GPU) + Q4_K_M (7B, CPU)",
     "양자화: AWQ 4-bit (32B) + Q4_K_M (7B)"),
    ("32B: GPU 고성능 분석 (H100/A100) + 7B: CPU 경량 분석",
     "32B GPU 고성능 + 7B CPU 경량 하이브리드"),
]

count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        for para in shape.text_frame.paragraphs:
            for old, new in TRIM:
                if old in para.text:
                    for run in para.runs:
                        if old in run.text:
                            run.text = run.text.replace(old, new)
                            count += 1
                        # multi-run case
                    full = ''.join(r.text for r in para.runs)
                    if old in full:
                        para.runs[0].text = full.replace(old, new)
                        for r in para.runs[1:]:
                            r.text = ""
                        count += 1

print(f"간결화: {count}건")
prs.save(DST)
print(f"저장: {DST}")
