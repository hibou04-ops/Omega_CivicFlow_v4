"""
Omega CivicFlow v4 — PPT 자동 수정 스크립트
기존 디자인/서식 보존하면서 텍스트만 업데이트
"""
import copy
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pathlib import Path

SRC = Path(r'c:\Users\hibou\Downloads\Omega_CivicFlow_v4_updated.pptx')
DST = Path(r'c:\Users\hibou\Downloads\Omega_CivicFlow_v4_final.pptx')

prs = Presentation(str(SRC))
slides = list(prs.slides)

def find_and_replace_in_slide(slide, old_text, new_text):
    """슬라이드 내 모든 shape에서 old_text를 찾아 new_text로 교체 (서식 보존)"""
    replaced = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full_text = para.text
            if old_text in full_text:
                # 단일 run에서 찾기
                for run in para.runs:
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)
                        replaced = True
                        break
                else:
                    # 다중 run에 걸쳐있는 경우 — 전체 run을 합쳐서 교체
                    if old_text in full_text and not replaced:
                        # run들의 텍스트를 합쳐서 교체 후 첫 run에 넣고 나머지 비움
                        if para.runs:
                            combined = full_text.replace(old_text, new_text)
                            para.runs[0].text = combined
                            for r in para.runs[1:]:
                                r.text = ""
                            replaced = True
    return replaced

def log_replace(slide_num, old, new, success):
    status = "✅" if success else "⚠️ NOT FOUND"
    print(f"  Slide {slide_num}: {status} '{old[:50]}...' → '{new[:50]}...'")

# ═══════════════════════════════════════════════════════
# 슬라이드 1 — 표지
# ═══════════════════════════════════════════════════════
s = slides[0]
old = "React 18  ·  FastAPI  ·  Qwen 2.5 32B  ·  Gemini 2.5 Pro  ·  ChromaDB  ·  Cognitive Engine"
new = "React 18  ·  FastAPI  ·  Qwen 2.5 32B  ·  Gemini 2.5 Pro/Flash  ·  PostgreSQL  ·  ChromaDB  ·  Cognitive Engine"
r = find_and_replace_in_slide(s, old, new)
log_replace(1, old, new, r)

# ═══════════════════════════════════════════════════════
# 슬라이드 2 — 목차 (09, 12 수정)
# ═══════════════════════════════════════════════════════
s = slides[1]

old = "Gemini 2.5 Pro — 전략 등급 S/A/B/C"
new = "Gemini Pro 분석 → Flash 감독 검증"
r = find_and_replace_in_slide(s, old, new)
log_replace(2, old, new, r)

old = "전략 Insight 엔진"
new = "전략 Insight + Supervisor"
r = find_and_replace_in_slide(s, old, new)
log_replace(2, old, new, r)

old = "SQLite 11 Tables"
new = "PostgreSQL 11 Tables + Supervisor 확장"
r = find_and_replace_in_slide(s, old, new)
log_replace(2, old, new, r)

# ═══════════════════════════════════════════════════════
# 슬라이드 6 — 핵심 차별점 (슬라이드 인덱스 5)
# ═══════════════════════════════════════════════════════
s = slides[5]

old = "전략 등급 시스템"
new = "2단 Insight 엔진"
r = find_and_replace_in_slide(s, old, new)
log_replace(6, old, new, r)

old = "Gemini 2.5 Pro → 투자 시사점·리스크·S/A/B/C 등급 자동 부여"
new = "Gemini Pro(분석) → Flash Supervisor(검증) — 사실/가정/미확인 분해 + 신뢰 등급 캘리브레이션"
r = find_and_replace_in_slide(s, old, new)
log_replace(6, old, new, r)

# ═══════════════════════════════════════════════════════
# 슬라이드 7 — 백엔드 서비스 레이어 (슬라이드 인덱스 6)
# ═══════════════════════════════════════════════════════
s = slides[6]

old = "21개 전문 서비스 모듈"
new = "24개 전문 서비스 모듈"
r = find_and_replace_in_slide(s, old, new)
log_replace(7, old, new, r)

old = "Gemini 2.5 Pro로 투자 시사점·리스크·전략 등급 생성 (암호화 프롬프트)"
new = "Gemini 2.5 Pro 1차 분석(암호화) + Flash Supervisor 사후 검증·보강"
r = find_and_replace_in_slide(s, old, new)
log_replace(7, old, new, r)

# ═══════════════════════════════════════════════════════
# 슬라이드 9 — LLM 분석 엔진 (슬라이드 인덱스 8)
# ═══════════════════════════════════════════════════════
s = slides[8]
old = "Ollama 로컬 실행 → API GPU 고성능 + CPU 경량 하이브리드"
new = "Ollama 로컬 실행 → API 비용 0원, GPU/CPU 하이브리드"
r = find_and_replace_in_slide(s, old, new)
log_replace(9, old, new, r)

# ═══════════════════════════════════════════════════════
# 슬라이드 11 — 전략 Insight 엔진 (슬라이드 인덱스 10) ★핵심★
# ═══════════════════════════════════════════════════════
s = slides[10]

# 제목 수정
old = "전략 Insight 엔진 (Gemini 2.5 Pro)"
new = "전략 Insight + Omega-Prime Supervisor"
r = find_and_replace_in_slide(s, old, new)
log_replace(11, old, new, r)

old = "The-Absolute — 투자 시사점 자동 생성"
new = "2단 구조: Gemini Pro 분석 → Flash 감독 검증"
r = find_and_replace_in_slide(s, old, new)
log_replace(11, old, new, r)

# 우측 "왜 Gemini 2.5 Pro인가?" 섹션 수정
old = "왜 Gemini 2.5 Pro인가?"
new = "2단계: Omega-Prime Supervisor (Flash)"
r = find_and_replace_in_slide(s, old, new)
log_replace(11, old, new, r)

old = "100만 토큰 컨텍스트 → 대량 문서 전체 분석 가능"
new = "1차 Insight를 5단계 추론 프로토콜로 사후 검증"
r = find_and_replace_in_slide(s, old, new)
log_replace(11, old, new, r)

old = "멀티모달 → 향후 차트/표 이미지 직접 분석 확장"
new = "사실/가정/미확인 분해 + 인과관계 검증"
r = find_and_replace_in_slide(s, old, new)
log_replace(11, old, new, r)

old = "한국어 성능 GPT-4o급 (2025 벤치마크 기준)"
new = "반론 생성(Counterfactual) + 신뢰 등급 캘리브레이션"
r = find_and_replace_in_slide(s, old, new)
log_replace(11, old, new, r)

old = "GCP 크레딧 활용 가능 → 초기 비용 절감"
new = "Pydantic v2 스키마 I/O 계약 + Flash로 비용 ~90% 절감"
r = find_and_replace_in_slide(s, old, new)
log_replace(11, old, new, r)

# ═══════════════════════════════════════════════════════
# 슬라이드 14 — DB 스키마 (슬라이드 인덱스 13)
# ═══════════════════════════════════════════════════════
s = slides[13]

old = "SQLite + SQLAlchemy ORM"
new = "PostgreSQL + SQLAlchemy ORM"
r = find_and_replace_in_slide(s, old, new)
log_replace(14, old, new, r)

# DocumentInsight 설명 업데이트
old = "Gemini 심층 분석"
new = "Gemini Pro + Flash Supervisor 통합 저장"
r = find_and_replace_in_slide(s, old, new)
log_replace(14, old, new, r)

# ═══════════════════════════════════════════════════════
# 슬라이드 16 — 기술 스택 (슬라이드 인덱스 15)
# ═══════════════════════════════════════════════════════
s = slides[15]

old = "SQLite + SQLAlchemy ORM"
new = "PostgreSQL + SQLAlchemy ORM"
r = find_and_replace_in_slide(s, old, new)
log_replace(16, old, new, r)

old = "설치 불필요, 단일 파일 배포"
new = "ACID 트랜잭션, 동시 접속 지원"
r = find_and_replace_in_slide(s, old, new)
log_replace(16, old, new, r)

old = "Gemini 2.5 Pro (GCP) × 2"
new = "Gemini 2.5 Pro + Flash (GCP)"
r = find_and_replace_in_slide(s, old, new)
log_replace(16, old, new, r)

old = "Insight용 + Chat용 별도 키 운영"
new = "Pro: Insight+Chat / Flash: Supervisor"
r = find_and_replace_in_slide(s, old, new)
log_replace(16, old, new, r)

# ═══════════════════════════════════════════════════════
# 슬라이드 17 — 아키텍처 (슬라이드 인덱스 16)
# ═══════════════════════════════════════════════════════
s = slides[16]

old = "21 Services (단일 책임 원칙)"
new = "24 Services (단일 책임 원칙)"
r = find_and_replace_in_slide(s, old, new)
log_replace(17, old, new, r)

# Data Layer
for shape in s.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if "SQLite + SQLAlchemy ORM" in para.text:
                for run in para.runs:
                    if "SQLite" in run.text:
                        run.text = run.text.replace("SQLite", "PostgreSQL")
                        print(f"  Slide 17: ✅ SQLite → PostgreSQL")
                        break

# 외부 서비스에 Flash 추가
old = "Gemini 2.5 Pro (GCP) — 전략 Insight + RAG 챗봇"
new = "Gemini 2.5 Pro/Flash (GCP) — Insight + Supervisor + 챗봇"
r = find_and_replace_in_slide(s, old, new)
log_replace(17, old, new, r)

# ═══════════════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════════════
prs.save(str(DST))
print(f"\n✅ 수정 완료! 저장 위치: {DST}")
